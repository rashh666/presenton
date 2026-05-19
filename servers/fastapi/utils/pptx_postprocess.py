"""
Post-processing for exported PPTX files driven by persona config.

Applies two optional decorations to every slide:
  1. A solid colour bar at the bottom edge (add_bottom_color_bar).
  2. A small watermark/logo image at the bottom-right corner (signature_watermark).

The function degrades gracefully: if python-pptx is absent or the file
cannot be opened/saved, a warning is logged and the original file is left
untouched.
"""

import logging
import os
from typing import Any, Dict, Optional

LOGGER = logging.getLogger(__name__)

_BAR_HEIGHT_EMU = 228600   # 0.24 inch
_WATERMARK_HEIGHT_EMU = 457200  # 0.5 inch
_WATERMARK_MARGIN_EMU = 114300  # 0.12 inch


def _hex_to_rgb(hex_colour: str):
    """Convert '#RRGGBB' to (R, G, B) integers."""
    h = hex_colour.lstrip("#")
    if len(h) != 6:
        raise ValueError(f"Invalid colour: {hex_colour!r}")
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def _resolve_watermark_path(raw: str, app_data_dir: str) -> Optional[str]:
    """
    Resolve an /app_data/... path to a real filesystem path.
    Returns None if the file does not exist.
    """
    if raw.startswith("/app_data/"):
        rel = raw[len("/app_data/"):]
        resolved = os.path.join(app_data_dir, rel)
    else:
        resolved = raw

    if os.path.isfile(resolved):
        return resolved
    LOGGER.debug("Watermark not found at %s — skipping", resolved)
    return None


def apply_persona_postprocess(
    pptx_path: str,
    persona_config: Dict[str, Any],
) -> None:
    """
    Modify *pptx_path* in-place according to persona post_processing settings.
    Does nothing (logs a warning) if python-pptx is unavailable or the config
    has no post_processing key.
    """
    post = persona_config.get("post_processing", {})
    if not post:
        return

    add_bar = post.get("add_bottom_color_bar", False)
    bar_colour_index = int(post.get("bottom_bar_color_index", 1))
    watermark_raw = post.get("signature_watermark", "")

    if not add_bar and not watermark_raw:
        return

    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Emu
    except ImportError:
        LOGGER.warning(
            "python-pptx is not installed; skipping PPTX post-processing for persona. "
            "Run: pip install python-pptx"
        )
        return

    palette = persona_config.get("visual_layout", {}).get("visual_palette", [])
    app_data_dir = (os.getenv("APP_DATA_DIRECTORY") or "/app_data").strip()

    try:
        prs = Presentation(pptx_path)
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        for slide in prs.slides:
            if add_bar and palette:
                _add_colour_bar(
                    slide,
                    palette,
                    bar_colour_index,
                    slide_width,
                    slide_height,
                )

            if watermark_raw:
                wm_path = _resolve_watermark_path(watermark_raw, app_data_dir)
                if wm_path:
                    _add_watermark(slide, wm_path, slide_width, slide_height)

        prs.save(pptx_path)
        LOGGER.info("Persona post-processing applied to %s", pptx_path)

    except Exception as exc:
        LOGGER.warning(
            "PPTX post-processing failed for %s: %s — original file preserved",
            pptx_path,
            exc,
        )


def _add_colour_bar(slide, palette, colour_index, slide_width, slide_height):
    from pptx.dml.color import RGBColor
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: F401 (used implicitly)

    idx = colour_index if colour_index < len(palette) else 0
    hex_colour = palette[idx]
    try:
        r, g, b = _hex_to_rgb(hex_colour)
    except ValueError as exc:
        LOGGER.debug("Bad colour %r: %s", hex_colour, exc)
        return

    left = Emu(0)
    top = Emu(slide_height - _BAR_HEIGHT_EMU)
    width = Emu(slide_width)
    height = Emu(_BAR_HEIGHT_EMU)

    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height,
    )
    shape.line.fill.background()  # no border
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)

    # Push bar behind other content
    sp_tree = slide.shapes._spTree
    sp_tree.remove(shape._element)
    sp_tree.insert(2, shape._element)


def _add_watermark(slide, image_path, slide_width, slide_height):
    from pptx.util import Emu

    height = Emu(_WATERMARK_HEIGHT_EMU)
    margin = Emu(_WATERMARK_MARGIN_EMU)

    try:
        from PIL import Image as _PIL_Image
        with _PIL_Image.open(image_path) as img:
            w_px, h_px = img.size
        aspect = w_px / h_px if h_px else 1.0
    except Exception:
        aspect = 1.0

    width = Emu(int(_WATERMARK_HEIGHT_EMU * aspect))
    left = Emu(slide_width - width - margin)
    top = Emu(slide_height - height - margin)

    slide.shapes.add_picture(image_path, left, top, width, height)
